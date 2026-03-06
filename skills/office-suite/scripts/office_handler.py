#!/usr/bin/env python3
"""
Office 三件套处理脚本 - 统一入口
支持 Word(.docx)、Excel(.xlsx)、PPT(.pptx) 的读取、创建和修改
"""

import sys
import os
import json
import argparse
from pathlib import Path

def handle_word(action, input_path=None, output_path=None, data=None):
    """处理 Word 文档"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return {"error": "请先安装 python-docx: pip install python-docx"}
    
    try:
        if action == "create":
            doc = Document()
            if data and "content" in data:
                for item in data["content"]:
                    if item.get("type") == "heading":
                        doc.add_heading(item.get("text", ""), level=item.get("level", 1))
                    elif item.get("type") == "paragraph":
                        p = doc.add_paragraph(item.get("text", ""))
                        if "style" in item:
                            p.style = item["style"]
                    elif item.get("type") == "table":
                        rows = item.get("rows", 2)
                        cols = item.get("cols", 2)
                        table = doc.add_table(rows=rows, cols=cols)
                        table.style = 'Light Grid Accent 1'
                        if "data" in item:
                            for i, row_data in enumerate(item["data"]):
                                if i < rows:
                                    row = table.rows[i]
                                    for j, cell_text in enumerate(row_data):
                                        if j < cols:
                                            row.cells[j].text = str(cell_text)
            doc.save(output_path)
            return {"success": True, "path": output_path, "action": "create"}
        
        elif action == "read":
            doc = Document(input_path)
            content = []
            full_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append({"type": "paragraph", "text": para.text})
                    full_text.append(para.text)
            
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                content.append({"type": "table", "data": table_data})
            
            return {
                "success": True,
                "path": input_path,
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "content": content,
                "full_text": "\n".join(full_text)
            }
        
        elif action == "modify":
            doc = Document(input_path)
            if data and "replacements" in data:
                for old_text, new_text in data["replacements"].items():
                    for para in doc.paragraphs:
                        if old_text in para.text:
                            para.text = para.text.replace(old_text, new_text)
            doc.save(output_path or input_path)
            return {"success": True, "path": output_path or input_path, "action": "modify"}
        
        else:
            return {"error": f"不支持的 Word 操作: {action}"}
    
    except Exception as e:
        return {"error": str(e)}


def handle_excel(action, input_path=None, output_path=None, data=None):
    """处理 Excel 表格"""
    try:
        import openpyxl
        from openpyxl import Workbook, load_workbook
    except ImportError:
        return {"error": "请先安装 openpyxl: pip install openpyxl"}
    
    try:
        if action == "create":
            wb = Workbook()
            ws = wb.active
            
            if data:
                if "sheet_name" in data:
                    ws.title = data["sheet_name"]
                if "headers" in data:
                    ws.append(data["headers"])
                if "rows" in data:
                    for row in data["rows"]:
                        ws.append(row)
            
            wb.save(output_path)
            return {"success": True, "path": output_path, "action": "create"}
        
        elif action == "read":
            wb = load_workbook(input_path, data_only=True)
            result = {"sheets": [], "active_sheet": wb.active.title}
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append(list(row))
                result["sheets"].append({
                    "name": sheet_name,
                    "rows": ws.max_row,
                    "cols": ws.max_column,
                    "data": rows
                })
            
            return {"success": True, "path": input_path, **result}
        
        elif action == "modify":
            wb = load_workbook(input_path)
            ws = wb.active
            
            if data:
                if "cell" in data:
                    cell_ref = data["cell"]
                    ws[cell_ref] = data.get("value", "")
                if "row" in data:
                    ws.append(data["row"])
            
            wb.save(output_path or input_path)
            return {"success": True, "path": output_path or input_path, "action": "modify"}
        
        else:
            return {"error": f"不支持的 Excel 操作: {action}"}
    
    except Exception as e:
        return {"error": str(e)}


def handle_ppt(action, input_path=None, output_path=None, data=None):
    """处理 PowerPoint 演示文稿"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"error": "请先安装 python-pptx: pip install python-pptx"}
    
    try:
        if action == "create":
            prs = Presentation()
            
            if data and "slides" in data:
                for slide_data in data["slides"]:
                    layout_idx = slide_data.get("layout", 1)
                    slide_layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts)-1)]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    if "title" in slide_data:
                        slide.shapes.title.text = slide_data["title"]
                    
                    if "content" in slide_data:
                        body_shape = slide.placeholders[1]
                        tf = body_shape.text_frame
                        tf.text = slide_data["content"]
            else:
                # 创建默认标题幻灯片
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = data.get("title", "标题") if data else "标题"
            
            prs.save(output_path)
            return {"success": True, "path": output_path, "action": "create"}
        
        elif action == "read":
            prs = Presentation(input_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_info = {"index": i + 1, "shapes": []}
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_info["shapes"].append({
                            "type": shape.shape_type,
                            "text": shape.text
                        })
                slides.append(slide_info)
            
            return {
                "success": True,
                "path": input_path,
                "slide_count": len(prs.slides),
                "slides": slides
            }
        
        elif action == "modify":
            prs = Presentation(input_path)
            
            if data and "add_slide" in data:
                layout_idx = data["add_slide"].get("layout", 1)
                slide_layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts)-1)]
                slide = prs.slides.add_slide(slide_layout)
                
                if "title" in data["add_slide"]:
                    slide.shapes.title.text = data["add_slide"]["title"]
            
            prs.save(output_path or input_path)
            return {"success": True, "path": output_path or input_path, "action": "modify"}
        
        else:
            return {"error": f"不支持的 PPT 操作: {action}"}
    
    except Exception as e:
        return {"error": str(e)}


def detect_file_type(path):
    """根据扩展名检测文件类型"""
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return "word"
    elif ext in [".xlsx", ".xlsm"]:
        return "excel"
    elif ext == ".pptx":
        return "ppt"
    else:
        return None


def main():
    parser = argparse.ArgumentParser(description="Office 三件套处理工具")
    parser.add_argument("--type", choices=["word", "excel", "ppt", "auto"], default="auto",
                       help="文件类型 (auto=自动检测)")
    parser.add_argument("--action", choices=["create", "read", "modify"], required=True,
                       help="操作类型")
    parser.add_argument("--input", "-i", help="输入文件路径")
    parser.add_argument("--output", "-o", help="输出文件路径")
    parser.add_argument("--data", "-d", help="JSON 格式的操作数据")
    
    args = parser.parse_args()
    
    # 自动检测文件类型
    file_type = args.type
    if file_type == "auto":
        target_path = args.input or args.output
        if target_path:
            file_type = detect_file_type(target_path)
        if not file_type:
            print(json.dumps({"error": "无法自动检测文件类型，请使用 --type 指定"}, ensure_ascii=False))
            sys.exit(1)
    
    # 解析数据
    data = None
    if args.data:
        try:
            data = json.loads(args.data)
        except json.JSONDecodeError as e:
            print(json.dumps({"error": f"JSON 解析错误: {e}"}, ensure_ascii=False))
            sys.exit(1)
    
    # 路由到对应的处理函数
    handlers = {
        "word": handle_word,
        "excel": handle_excel,
        "ppt": handle_ppt
    }
    
    handler = handlers.get(file_type)
    if not handler:
        print(json.dumps({"error": f"不支持的文件类型: {file_type}"}, ensure_ascii=False))
        sys.exit(1)
    
    result = handler(args.action, args.input, args.output, data)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
